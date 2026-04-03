#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PharmPay IR Pitchdeck Generator
최종 검증된 데이터 기반 투자자 미팅용 피치덱
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_pharmpay_pitchdeck():
    """PharmPay IR 피치덱 생성"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 색상 팔레트 (전문적인 블루 계열)
    PRIMARY_BLUE = RGBColor(0, 102, 204)      # #0066CC
    SECONDARY_BLUE = RGBColor(51, 153, 255)   # #3399FF
    ACCENT_GREEN = RGBColor(0, 153, 76)       # #00994C
    DARK_GRAY = RGBColor(51, 51, 51)          # #333333
    LIGHT_GRAY = RGBColor(242, 242, 242)      # #F2F2F2
    
    def add_title_slide(title, subtitle=""):
        """타이틀 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # 배경색
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_GRAY
        
        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # 서브타이틀
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = DARK_GRAY
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_list, footer=""):
        """컨텐츠 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE
        
        # 구분선
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0.5), Inches(1.2), Inches(9), Inches(0)
        )
        line.line.color.rgb = SECONDARY_BLUE
        line.line.width = Pt(2)
        
        # 컨텐츠
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.2)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        for i, item in enumerate(content_list):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)
            p.level = 0 if not item.startswith("  ") else 1
        
        # Footer
        if footer:
            footer_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(7.0), Inches(9), Inches(0.4)
            )
            footer_frame = footer_box.text_frame
            footer_frame.text = footer
            footer_para = footer_frame.paragraphs[0]
            footer_para.font.size = Pt(10)
            footer_para.font.color.rgb = DARK_GRAY
            footer_para.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    # ==================== Slide 1: 커버 ====================
    add_title_slide(
        "PharmPay",
        "처방전 기반 독점 RWD 데이터를 활용한\n약국 AI 복약지도 플랫폼"
    )
    
    # ==================== Slide 2: Problem ====================
    add_content_slide(
        "Problem: 약국의 3가지 고통",
        [
            "1️⃣ 약사 비효율",
            "  • 복약지도 10~15분 소요, 하루 30~50건 반복",
            "  • 번아웃 심각, 10곳 중 4곳 약사 구하기 어려움",
            "",
            "2️⃣ 환자 불만",
            "  • 약국 대기 시간 길고, 의학 용어 이해 어려움",
            "  • 집에 가면 복약 방법 잊어버림 (복약 순응도 50%)",
            "",
            "3️⃣ 수익성 악화",
            "  • 복약지도는 비용 (무료), 약가 인하 압력",
            "  • 온라인 약국 경쟁, 약국 평균 영업이익률 7.8% → 6.5%"
        ],
        "출처: 대한약사회 (2024), 한국소비자원 (2025)"
    )
    
    # ==================== Slide 3: Solution ====================
    add_content_slide(
        "Solution: PharmPay AI 플랫폼",
        [
            "📱 처방전 스캔 → AI 분석 → 맞춤 복약 가이드 생성 → 약사 검토 → 환자 전송",
            "",
            "핵심 기능:",
            "  • AI 복약지도 자동화: Claude/Gemini 기반, 약물 상호작용 체크",
            "  • 환자 맞춤 가이드: 쉬운 언어로 복약 타이밍 설명 (알림톡)",
            "  • 건기식 추천: 처방약 기반 맞춤 영양제 (약국 매출 증대)",
            "  • 단골 관리: 환자 복약 이력 관리, 재구매 알림",
            "",
            "✅ 결과:",
            "  • 복약지도 시간 50% 단축 (10분 → 5분)",
            "  • 약사 만족도 4.5/5.0 (PoC 10개소)",
            "  • 약국 OTC/건기식 매출 15% 증가"
        ],
        "PoC 검증 완료 (2025 Q4 ~ 2026 Q1)"
    )
    
    # ==================== Slide 4: Traction ====================
    add_content_slide(
        "Traction: 검증된 성과",
        [
            "🎯 PoC 완료 (10개 약국 운영 중)",
            "  • 약사 만족도: 4.5/5.0",
            "  • 환자 긍정 피드백: 92%",
            "  • 복약 시간 단축: 50%",
            "",
            "🤝 크레소티 파트너십 체결 (2026 Q1)",
            "  • 약국 POS 업계 1위 (시장점유율 약 60%, 1.4만개 약국)",
            "  • 공동 마케팅 합의: 크레소티 POS에 PharmPay 번들링",
            "  • CAC 80% 절감: 직접 영업 100만원 → 번들링 20만원",
            "",
            "💻 기술 검증",
            "  • OCR 정확도: 98% (처방전 스캔)",
            "  • AI 응답 시간: 평균 3초",
            "  • 약물 DB: 2만개+ 구축 완료"
        ],
        "크레소티: 업계 추정 약 60% 점유율"
    )
    
    # ==================== Slide 5: Market Size ====================
    add_content_slide(
        "Market Opportunity: 15조원 시장",
        [
            "📊 TAM (Total Addressable Market): 15조원",
            "  • OTC 의약품: 8조원 (식약처 2023)",
            "  • 건강기능식품: 6조원 (한국건기식협회 2025)",
            "  • 디지털 헬스케어: 1조원 (추정)",
            "",
            "🎯 SAM (Serviceable Available Market): 2,400억원",
            "  • 약국 SaaS 시장: 약 2.4만 약국 × 연 1,000만원",
            "",
            "💰 SOM (Serviceable Obtainable Market): 500억원",
            "  • 2028년 목표: 3,000개 약국 (12.5%)",
            "  • 연 매출: 100억원 (보수적 추정)",
            "",
            "🏆 경쟁 환경: 직접 경쟁자 없음 (Blue Ocean)",
            "  • 메디히어: 원격진료 (진료 플랫폼, 다른 영역)",
            "  • 메딜리티: 알약 카운팅 (업무 자동화, 보완재)"
        ],
        "출처: 식약처(2023), 한국건기식협회(2025)"
    )
    
    # ==================== Slide 6: Business Model ====================
    add_content_slide(
        "Business Model: 3단 로켓 수익 모델",
        [
            "🚀 1단: SaaS 구독 (Year 1~2)",
            "  • 약국당 월 15~50만원 (Basic → Pro → Premium)",
            "  • 안정적 Recurring Revenue",
            "",
            "🚀 2단: Ad-Tech 광고 (Year 2~3)",
            "  • 제약사 타겟 광고 (고지혈증 약 → 코큐텐 광고)",
            "  • 분배: 약국 40% / PharmPay 60%",
            "  • CPM: 5~10만원 (일반 포털 대비 3~5배 프리미엄)",
            "",
            "🚀 3단: RWD 데이터 판매 (Year 3+)",
            "  • 제약사 라이선스: 연 5,000만원/사",
            "  • 용도: 신약 PMS, 약물 안전성 모니터링, 마케팅",
            "  • 마진: 95%+ (Pure Data Play)",
            "",
            "💎 독점 데이터 = 국내 유일 처방전 기반 RWD"
        ],
        "식약처 PMS(시판 후 조사) 의무화로 RWD 수요 급증"
    )
    
    # ==================== Slide 7: Unit Economics ====================
    add_content_slide(
        "Unit Economics: 건강한 SaaS 모델",
        [
            "📈 핵심 지표 (보수적 시나리오)",
            "  • LTV/CAC: 28.8배 (목표 3배 대비 9배 우수) ✅",
            "  • CAC Payback: 2개월 (목표 6개월 대비 3배 빠름) ✅",
            "  • Gross Margin: 75%+ (고마진 SaaS) ✅",
            "  • Churn Rate: 5~10% (낮은 이탈률) ✅",
            "",
            "📊 약국당 경제성",
            "  • ARPU: 월 15만원 → 100만원 (3년)",
            "  • CAC: 20~30만원 (크레소티 번들링)",
            "  • LTV: 864만원 (보수적)",
            "",
            "💰 손익분기점",
            "  • BEP: 2027년 Q2 (200개 약국, 월 매출 1억원)",
            "  • 필요 약국: 200개 (ARPU 50만원 기준)"
        ],
        "Claude + Gemini 교차 검증 완료"
    )
    
    # ==================== Slide 8: Financial Projections ====================
    add_content_slide(
        "Financial Projections: 3개년 전망 (보수적)",
        [
            "📅 2026년",
            "  • 연동 약국: 100개",
            "  • 연 매출: 0.3억원",
            "  • 상태: PoC 확장, 제품 고도화",
            "",
            "📅 2027년 (권고안: 상향 조정)",
            "  • 연동 약국: 1,000개",
            "  • MAU: 10만명",
            "  • 연 매출: 30억원 (기존 5억원 → 상향)",
            "  • 상태: BEP 달성 (Q2)",
            "",
            "📅 2028년 (권고안: 상향 조정)",
            "  • 연동 약국: 3,000개 (시장 12.5%)",
            "  • MAU: 100만명",
            "  • 연 매출: 100억원 (기존 20억원 → 상향)",
            "  • 상태: Series A 완료, RWD 수익화 시작",
            "",
            "💡 실제 가능 매출: 2027년 48억원, 2028년 252억원",
            "    (Unit Economics 기준 계산)"
        ],
        "보수적 추정: 실제 가능 매출의 40~60% 수준"
    )
    
    # ==================== Slide 9: Competitive Advantage ====================
    add_content_slide(
        "Competitive Advantage: 4가지 차별점",
        [
            "1️⃣ 독점 데이터 자산 (핵심 경쟁력)",
            "  • 처방전 기반 RWD = 국내 유일",
            "  • 제약사/보험사에게 금맥 (PMS 의무화)",
            "  • 후발 주자는 0부터 시작 (데이터 해자)",
            "",
            "2️⃣ 크레소티 독점 파트너십 (GTM 확보)",
            "  • 2.4만 약국 중 60% 즉시 접근 가능",
            "  • 경쟁사 대비 10배 빠른 시장 진입",
            "",
            "3️⃣ 건강한 Unit Economics",
            "  • LTV/CAC 28배 (SaaS 최고 수준)",
            "  • CAC Payback 2개월 (빠른 현금화)",
            "",
            "4️⃣ 규제 친화적",
            "  • 의료기기 아님 (식약처 인허가 불필요)",
            "  • 약사법 준수 (AI는 보조, 약사가 최종 판단)",
            "  • PIMS 인증 준비 중 (개인정보보호)"
        ],
        "법무법인 의견서 확보 예정"
    )
    
    # ==================== Slide 10: Team ====================
    add_content_slide(
        "Team: 하드웨어/소프트웨어 융합 + AI 기술력",
        [
            "👨‍💼 김상현 (CEO) - 10년차 헬스케어 서비스 상용화 전문가",
            "  • UIUC (일리노이대학교) 산업공학",
            "  • IoT 헬스케어 디바이스 4만 대 양산 경험",
            "    - (주)디디케어스: 반려동물 웨어러블 3만 대",
            "    - (주)아이디엘: 아동 IoT 디바이스 1만 대",
            "  • 하드웨어/플랫폼 융합 End-to-End 실행력",
            "",
            "👨‍💻 김종우 (CTO) - 200만 MAU 운영 경험의 AI 아키텍트",
            "  • UCLA 컴퓨터공학, AI/NLP 특허 14건",
            "  • 현대차/아모레 AI 챗봇 상용화 (2년 운영)",
            "  • LLM 기반 서비스 5건 출시 (RAG 기술 확보)",
            "  • 글로벌 웹서비스 MAU 200만 명 무중단 운영",
            "  • Web3 보안 기술 (의료 데이터 암호화)",
            "",
            "🎯 Advisory Board: 약사회, 헬스케어 법률, 의료 AI 전문가"
        ],
        "POS(하드웨어) + AI(소프트웨어) 융합 경험 완비"
    )
    
    # ==================== Slide 11: Why Now? ====================
    add_content_slide(
        "Why Now? 2026년이 최적 타이밍",
        [
            "1️⃣ 기술 성숙",
            "  • LLM 혁명으로 복약지도 자동화 가능 (Claude/Gemini)",
            "  • AI 의료 벤치마크 통과 (약사 수준 90%)",
            "",
            "2️⃣ 시장 준비",
            "  • 약국 60%가 POS 디지털화 완료 (크레소티 등)",
            "  • 약사 자동화 수요 폭발 (10곳 중 4곳 인력난)",
            "",
            "3️⃣ 정책 지원",
            "  • 정부 디지털 헬스케어 육성 정책 (연 500억원)",
            "  • 규제 샌드박스 + R&D 예산 확대",
            "  • 건강보험 디지털 청구 의무화 (2025년)",
            "",
            "4️⃣ 경쟁 공백",
            "  • 약국 AI 복약지도 직접 경쟁자 없음 (Blue Ocean)",
            "",
            "5️⃣ 데이터 기회",
            "  • 제약사 RWD 수요 급증 (PMS 의무화)",
            "  • 처방전 데이터 = 독점 자산"
        ],
        "지금 진입하지 않으면 기회 상실"
    )
    
    # ==================== Slide 12: Investment Ask ====================
    add_content_slide(
        "Investment Ask: Seed Round 3~5억원",
        [
            "💰 투자 조건",
            "  • 조달 금액: 3~5억원",
            "  • Pre-money 밸류에이션: 18억원",
            "  • Post-money 밸류에이션: 22억원 (4억원 조달 기준)",
            "  • 투자자 희석률: 18~20%",
            "  • 창업자 지분 유지: 80%+",
            "",
            "🎯 Use of Funds (4억원 기준)",
            "  • 개발 (60%, 2.4억원): AI 고도화, 개발자 3명 채용",
            "  • 마케팅 (20%, 0.8억원): 약국 온보딩 100개소",
            "  • 인프라 (20%, 0.8억원): AWS, 보안, PIMS 인증",
            "",
            "📅 마일스톤",
            "  • M1 (6개월): 50개 약국",
            "  • M2 (12개월): 100개 약국, BEP 근접",
            "  • M3 (18개월): 500개 약국, Series A 준비"
        ],
        "24개월 런웨이 (BEP까지 충분)"
    )
    
    # ==================== Slide 13: Why Invest? ====================
    add_content_slide(
        "Why Invest in PharmPay?",
        [
            "✅ 1. 검증된 Product-Market Fit",
            "  • PoC 10개소 완료 (만족도 4.5/5.0)",
            "  • 크레소티 파트너십 (1.4만 약국 접근)",
            "  • 약사 ROI 명확 (복약 시간 50% 단축)",
            "",
            "✅ 2. 독점 데이터 자산",
            "  • 처방전 RWD = 국내 유일",
            "  • 제약사 B2B 수익화 (연 5억원+ 잠재력)",
            "  • 데이터 네트워크 효과 (선점 효과)",
            "",
            "✅ 3. 건강한 Unit Economics",
            "  • LTV/CAC 28배 (SaaS 최고 수준)",
            "  • CAC Payback 2개월 (빠른 현금화)",
            "  • Gross Margin 75%+ (고수익)",
            "",
            "✅ 4. Blue Ocean 시장",
            "  • 직접 경쟁자 없음, TAM 15조원",
            "",
            "✅ 5. 명확한 Exit 경로",
            "  • M&A (3~4년): 크레소티, 네이버/카카오, 대형 제약사",
            "  • Exit Value: 300~500억원 예상",
            "  • 투자자 ROI: 15~20배 (3년 기준)"
        ]
    )
    
    # ==================== Slide 14: Contact ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Let's Build the Future\nof Pharmacy Together"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.0), Inches(9), Inches(2)
    )
    contact_frame = contact_box.text_frame
    contact_text = """(주)페보 PharmPay
CEO: 김상현
Email: ceo@pharmpay.kr
Phone: 010-XXXX-XXXX

감사합니다."""
    contact_frame.text = contact_text
    for para in contact_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(6)
    
    # 파일 저장
    output_path = "/Users/sanghyunkim/Desktop/ai-company/IR Deck/PharmPay_IR_Pitchdeck_Final.pptx"
    prs.save(output_path)
    print(f"✅ 피치덱 생성 완료: {output_path}")
    print(f"📊 총 {len(prs.slides)} 슬라이드")
    
    return output_path

if __name__ == "__main__":
    create_pharmpay_pitchdeck()
