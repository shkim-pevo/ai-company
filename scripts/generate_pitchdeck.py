#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PharmPay IR Pitch Deck Generator
(주)페보 | Seed Round 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_pharmpay_pitchdeck():
    """PharmPay IR 피치덱 생성"""
    
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 브랜드 컬러
    BRAND_BLUE = RGBColor(28, 69, 135)  # #1C4587
    BRAND_TEAL = RGBColor(68, 114, 196)  # #4472C4
    ACCENT_GREEN = RGBColor(112, 173, 71)  # #70AD47
    DARK_GRAY = RGBColor(51, 51, 51)  # #333333
    LIGHT_GRAY = RGBColor(242, 242, 242)  # #F2F2F2
    WHITE = RGBColor(255, 255, 255)
    
    def add_title_slide(title, subtitle, phase=""):
        """타이틀 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # 배경색
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BRAND_BLUE
        
        # 메인 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # 서브타이틀
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = ACCENT_GREEN
        
        # Phase 표시
        if phase:
            phase_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
            phase_frame = phase_box.text_frame
            phase_frame.text = phase
            phase_para = phase_frame.paragraphs[0]
            phase_para.alignment = PP_ALIGN.CENTER
            phase_para.font.size = Pt(18)
            phase_para.font.color.rgb = LIGHT_GRAY
        
        return slide
    
    def add_content_slide(title, content_items, phase="", highlight=None):
        """컨텐츠 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 배경
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # 상단 Phase 바
        if phase:
            phase_bar = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0), Inches(0), Inches(10), Inches(0.5)
            )
            phase_bar.fill.solid()
            phase_bar.fill.fore_color.rgb = BRAND_BLUE
            phase_bar.line.fill.background()
            
            phase_text = phase_bar.text_frame
            phase_text.text = phase
            phase_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            phase_text.paragraphs[0].font.size = Pt(14)
            phase_text.paragraphs[0].font.color.rgb = WHITE
            phase_text.paragraphs[0].font.bold = True
        
        # 타이틀
        title_top = Inches(0.8) if phase else Inches(0.5)
        title_box = slide.shapes.add_textbox(Inches(0.5), title_top, Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = BRAND_BLUE
        
        # 하이라이트 박스 (있는 경우)
        content_top = Inches(1.8) if phase else Inches(1.5)
        if highlight:
            highlight_box = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), content_top, Inches(9), Inches(0.8)
            )
            highlight_box.fill.solid()
            highlight_box.fill.fore_color.rgb = LIGHT_GRAY
            highlight_box.line.color.rgb = BRAND_TEAL
            highlight_box.line.width = Pt(2)
            
            highlight_text = highlight_box.text_frame
            highlight_text.text = highlight
            highlight_text.paragraphs[0].font.size = Pt(18)
            highlight_text.paragraphs[0].font.color.rgb = DARK_GRAY
            highlight_text.paragraphs[0].font.bold = True
            highlight_text.margin_left = Inches(0.2)
            highlight_text.margin_top = Inches(0.15)
            
            content_top = Inches(2.8) if phase else Inches(2.5)
        
        # 컨텐츠 항목들
        for i, item in enumerate(content_items):
            item_box = slide.shapes.add_textbox(
                Inches(0.8), 
                content_top + Inches(i * 0.9), 
                Inches(8.4), 
                Inches(0.8)
            )
            item_frame = item_box.text_frame
            item_frame.word_wrap = True
            
            # 타이틀과 내용 분리
            if isinstance(item, dict):
                # 딕셔너리 형태: {'title': '...', 'content': '...'}
                p = item_frame.paragraphs[0]
                p.text = f"• {item['title']}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = BRAND_TEAL
                
                if item.get('content'):
                    p2 = item_frame.add_paragraph()
                    p2.text = item['content']
                    p2.font.size = Pt(14)
                    p2.font.color.rgb = DARK_GRAY
                    p2.space_before = Pt(6)
                    p2.level = 1
            else:
                # 문자열 형태
                p = item_frame.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
        
        return slide
    
    # ==================== Slide 1: 타이틀 ====================
    add_title_slide(
        "PharmPay",
        "내 손안의 나만의 AI 약사",
        ""
    )
    
    # 슬라이드 1 추가 정보
    slide1 = prs.slides[-1]
    info_box = slide1.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.text = "1.5만 약국 POS 연동 기반\n24시간 AI 복약케어 플랫폼\n\n(주)페보 | Seed Round 2026\ntech@pevo.care"
    for para in info_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE
    
    # ==================== Slide 2: 약사의 본질 ====================
    add_content_slide(
        "약사의 본질 (The Essence)",
        [
            {"title": "처방의 끝이 아닌 '케어의 시작'", 
             "content": "약사는 의료진 중 가장 접근성 높은 약물 전문 게이트키퍼"},
            {"title": "핵심 역할", 
             "content": "처방 검토 및 중재 (약물상호작용·알레르기 차단) + 전주기 복약 모니터링"},
            {"title": "케어 갭 현실", 
             "content": "일 평균 100건+ 처방 처리 → 환자당 상담시간 2분 미만 → 개인화 불가능"},
            {"title": "경제적 임팩트", 
             "content": "복약순응도 1% 향상 시 국가 의료비 절감 약 3,000억원 (HIRA 추정)"}
        ],
        "PHASE 1 — PROBLEM",
        "처방 검토와 전주기 복약 모니터링이 약사의 진짜 역할"
    )
    
    # ==================== Slide 3: 문제 정의 ====================
    add_content_slide(
        "문제 정의 (Problem)",
        [
            {"title": "환자의 불안 (Care Gap)", 
             "content": "무자문 구매율 67% — 밤중 부작용 발생 시 물어볼 채널 없어 인터넷 거짓정보 의존"},
            {"title": "약사의 한계 (물리적 단절)", 
             "content": "일 평균 100건+ 처리로 귀가 환자 추적·연락 방법 전무. 약사 번아웃율 62%"},
            {"title": "수익 기회 상실 (골든타임)", 
             "content": "단골 환자 유대감 끊어지며 OTC+건기식 8조원+ 시장 기회 상실"}
        ],
        "PHASE 1 — PROBLEM",
        "약국 문을 나서는 순간 '의료 사각지대'가 시작된다"
    )
    
    # ==================== Slide 4: 솔루션 ====================
    add_content_slide(
        "솔루션 (Solution)",
        [
            {"title": "Voice AI (24시간 밀착 케어)", 
             "content": "처방 연동 기반 카카오 알림톡으로 '오늘 불편한 곳은 없으셨나요?' 음성 대화"},
            {"title": "AI Recommendation (맞춤 큐레이션)", 
             "content": "처방 이력 기반 체내 결핍 영양소 분석, OTC·건기식 최적 추천 및 안전성 검증"},
            {"title": "O2O Care Link (약국 재연결)", 
             "content": "이상 징후 감지 시 약사 대시보드 즉시 알림 → 환자의 약국 재방문 유도"}
        ],
        "PHASE 2 — SOLUTION",
        "끊어진 선을 잇다 — 24시간 나를 기억하는 '내 손안의 AI 약사'"
    )
    
    # ==================== Slide 5: 압도적 해자 ====================
    add_content_slide(
        "압도적 해자 (Unfair Advantage)",
        [
            {"title": "크레소티 파트너십", 
             "content": "국내 약국 POS 점유율 75% 기업과 협업, 클릭 한 번으로 전국 배포"},
            {"title": "처방 데이터 자동 동기화", 
             "content": "결제 순간 처방 데이터가 AI와 연동, 환자 입력 필요 없음"},
            {"title": "선점 효과 (Data Moat)", 
             "content": "가장 넓은 '약사-환자 연결망' 선점, 후발 주자 진입 불가"}
        ],
        "PHASE 3 — BUSINESS MODEL",
        "국내 1위 POS 인프라 독점 연동 — 크레소티 POS 1.5만 약국, 점유율 75%"
    )
    
    # ==================== Slide 6: 시장 기회 ====================
    add_content_slide(
        "시장 기회 (Market Size)",
        [
            {"title": "TAM ~15조원", 
             "content": "OTC + 건기식 + 디지털약국 전체 시장"},
            {"title": "SAM ~2조원", 
             "content": "약국 채널 AI/디지털 전환 시장"},
            {"title": "SOM 500억원", 
             "content": "크레소티 연동 약국 기반 초기 시장"},
            {"title": "핵심 지표", 
             "content": "전국 약국 2.4만개소, OTC+건기식 8조원+, 만성질환자 1,700만명"}
        ],
        "PHASE 3 — BUSINESS MODEL",
        "약국 기반 디지털 헬스케어 시장"
    )
    
    # ==================== Slide 7: 수익 모델 ====================
    add_content_slide(
        "수익 모델 (Business Model)",
        [
            {"title": "SaaS 구독 (기저 매출, 2026 Q3~)", 
             "content": "약국 단골관리 대시보드, 복약지도 연동, 약사 알림 대시보드"},
            {"title": "Ad-Tech 타겟 광고 (성장 엔진, 2027 Q1~)", 
             "content": "실제 질환/복약 기반 타겟팅, 포털 대비 압도적 전환율, 고단가 CPM"},
            {"title": "B2B 데이터 판매 (Cash Cow, 2027 H2~)", 
             "content": "제약사 실제 복용 데이터(RWD), 임상/부작용(PMS) 리포트"}
        ],
        "PHASE 3 — BUSINESS MODEL",
        "초정밀 타겟 Ad-Tech & Data 플랫폼 — SaaS 기저매출 + 광고 + RWD 데이터 3축 수익"
    )
    
    # ==================== Slide 8: Go-To-Market 전략 ====================
    add_content_slide(
        "Go-To-Market 전략",
        [
            {"title": "Phase 1 (2026 H1) 실증(PoC)", 
             "content": "거점 우수 약국 10개소, 환자 재방문율 증명, 객단가 상승 데이터 확보"},
            {"title": "Phase 2 (2026 H2) 네트워크 확보", 
             "content": "팜페이 번들 프로모션, 유료 전환 약국 1,000개, 환자 네트워크 10만명"},
            {"title": "Phase 3 (2027~) 광고 매체화", 
             "content": "B2C 환자 트래픽 100만명, 제약사 마케팅 예산 유치, 월 매출 5억원 목표"}
        ],
        "PHASE 3 — BUSINESS MODEL",
        "크레소티 번들링을 통한 단기간 시장 장악 — 3단계 로드맵"
    )
    
    # ==================== Slide 9: 경쟁 환경 ====================
    add_content_slide(
        "경쟁 환경 (Competition)",
        [
            {"title": "크레소티 POS 독점 연동", 
             "content": "경쟁사 대비 10배 빠른 배포 (1.5만 약국 즉시 접근)"},
            {"title": "유일한 '처방 데이터 기반' 초개인화", 
             "content": "다른 앱은 사용자 자가 입력에 의존"},
            {"title": "약국을 '건강 허브'로 전환", 
             "content": "단순 판매가 아닌 사후 케어 수익화 구조"}
        ],
        "PHASE 3 — BUSINESS MODEL",
        "약국관리 SaaS, 헬스케어 앱, 배송 플랫폼 대비 유일한 '처방 데이터 기반' 통합 플랫폼"
    )
    
    # ==================== Slide 10: 재무 전망 ====================
    add_content_slide(
        "재무 전망 (Financials)",
        [
            {"title": "2026년", 
             "content": "연동 약국 10개소, 연 매출 ~0.3억원 (PoC 단계)"},
            {"title": "2027년", 
             "content": "연동 약국 1,000개, 활성 환자(MAU) 10만, ARPU 7만원/월, 연 매출 5억원"},
            {"title": "2028년", 
             "content": "연동 약국 3,000개, MAU 100만, ARPU 10만원/월, 연 매출 20억원, BEP 달성"}
        ],
        "PHASE 3 — BUSINESS MODEL",
        "3개년 보수적 전망 — 2028년 연 매출 20억원, BEP 달성 목표"
    )
    
    # ==================== Slide 11: Vision & Expansion ====================
    add_content_slide(
        "Vision & Expansion (궁극적 비전)",
        [
            {"title": "NOW", 
             "content": "복약 상담 AI + OTC·건기식 추천 (MVP 단계)"},
            {"title": "NEXT", 
             "content": "만성질환 관리 + 소아 케어 확장 — 생애주기별 맞춤 AI 시나리오 고도화"},
            {"title": "FUTURE", 
             "content": "반려동물 케어 + 가족 통합 구독 — 멀리 계신 부모님 복약 순응도를 자녀가 확인하는 가족 연동망"},
            {"title": "CEO Track Record", 
             "content": "아동 웨어러블 1만대 → 펫 헬스케어 3만대 → AI 약사 — 10년간 전 생애 헬스케어 빌드업"}
        ],
        "PHASE 4 — VISION & TEAM",
        "질병 관리를 넘어 '전 생애주기 맞춤형 가족 헬스케어'로"
    )
    
    # ==================== Slide 12: 팀 ====================
    add_content_slide(
        "팀 (Team)",
        [
            {"title": "CEO 김상현 (UIUC 산업공학)", 
             "content": "아동 웨어러블 1만대 양산, 펫 헬스케어 3만대 플랫폼 구축, End-to-End 생태계 경험"},
            {"title": "CTO 김종우 (UCLA 컴퓨터공학)", 
             "content": "1세대 AI 챗봇 개발(현대차/아모레), AI 특허 14건, 200만 MAU 무중단 운영"},
            {"title": "핵심 파트너", 
             "content": "크레소티(Cresoti) — 약국 POS 1위, AI 공동개발 파트너 + 약학 자문위원 위촉 예정"}
        ],
        "PHASE 4 — VISION & TEAM",
        "헬스케어의 점들을 연결해 온 연쇄 창업가 & AI 아키텍트"
    )
    
    # ==================== Slide 13: 투자 요청 ====================
    add_content_slide(
        "투자 요청 (The Ask)",
        [
            {"title": "시드 라운드 3~5억원", 
             "content": "12개월 내 약국 100개, 환자 10만, MRR 3,000만원 목표"},
            {"title": "자금 사용처", 
             "content": "핵심 인력(CTO/개발) 40%, AI 엔진/API 연동 30%, 파일럿/마케팅 20%, 운영/법무 10%"},
            {"title": "12개월 목표", 
             "content": "연동 약국 10 → 100개, 활성 환자 0 → 10만명, MRR 0 → 3,000만원"},
            {"title": "연락처", 
             "content": "(주)페보 | 대표 김상현 | tech@pevo.care | 경기도 성남시 | 설립 2024.03"}
        ],
        "PHASE 4 — VISION & TEAM",
        ""
    )
    
    # ==================== Slide 14: Closing ====================
    closing_slide = add_title_slide(
        "Thank You",
        "내 손안의 나만의 AI 약사",
        ""
    )
    
    # Contact 정보 추가
    contact_box = closing_slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "(주)페보\n대표 김상현\ntech@pevo.care\n\nPartner: 크레소티(Cresoti)"
    for para in contact_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE
    
    # 저장
    output_path = "/Users/sanghyunkim/Desktop/ai-company/IR Deck/PharmPay_IR_Pitchdeck_2026_v2.pptx"
    prs.save(output_path)
    print(f"✅ 피치덱이 성공적으로 생성되었습니다: {output_path}")
    print(f"📊 총 {len(prs.slides)}개의 슬라이드가 생성되었습니다.")
    
    return output_path

if __name__ == "__main__":
    create_pharmpay_pitchdeck()
